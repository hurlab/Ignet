#!/usr/bin/env python3
"""Revise the hand-edited ICIBM deck v3.1 -> v3.2.

This edits the user's deck IN PLACE (copy-then-modify). It deliberately does NOT
regenerate from build_icibm_deck_v21.py: v3.1 carries substantial manual work
(new use-case slides, screenshots, layout) that a rebuild would discard.

Changes, each independently verifiable:

  1. Strip every speaker note (user request: remove all).
  2. Times New Roman -> Trebuchet MS on the two use-case slides. Those were the
     only inconsistent slides -- the AI/ontology slide was already Trebuchet.
  3. "Aanalyze Your Own Text" -> "Analyze Your Own Text".
  4. Use case 2 reframed to lead with the question it answers, instead of
     describing the capability chain.
  5. NEW use case 3: the R01 preliminary-data analysis where Ignet's disease
     ontology layer supplied one of two literature-derived gene sets.
  6. The AI/ontology slide rebuilt as terms-vs-structure. The old version
     claimed the AI layer "runs on" the ontology layer; only three of the
     integrated ontologies are actually traversed. It also had no slide number,
     because it used a layout that carries no such placeholder.
  7. Conclusion trimmed from full sentences to phrases.
  8. Platform Benchmarking moved behind the closing slide as a backup.

Slide numbers are auto-fields (<a:fld type="slidenum">), so reordering and
insertion renumber themselves -- no manual fixups.

Usage:
    python3 scripts/revise_icibm_v32.py --in <v3.1.pptx> --out <v3.2.pptx>
"""
from __future__ import annotations

import argparse
import pathlib
import shutil
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Inches, Pt

BODY_FONT = "Trebuchet MS"
GREEN_BOX = RGBColor(0xE0, 0xEE, 0xD9)   # sampled from the use-case question box
CYAN_BOX = RGBColor(0xF0, 0xFE, 0xFB)    # sampled from the use-case caption box
TEAL = RGBColor(0x0E, 0x7C, 0x7B)        # its border
MINT = RGBColor(0xE6, 0xF5, 0xEC)        # sampled from the conclusion takeaway box
INK = RGBColor(0x2C, 0x3C, 0x43)

FIGURE = pathlib.Path(
    "/home/juhur/PROJECTS/60_grants/grant_rutgers/phase1/figs/v2b_fig5_themes.png"
)

log = []


def note(msg):
    log.append(msg)
    print(msg)


# --------------------------------------------------------------------------
# small helpers
# --------------------------------------------------------------------------

def title_of(slide):
    """Title with soft line breaks and non-breaking spaces normalised.

    At least one hand-edited title carries U+00A0 ("Platform\xa0Benchmarking"),
    which silently defeats a plain substring lookup.
    """
    if slide.shapes.title is None:
        return ""
    return (slide.shapes.title.text
            .replace("\x0b", " ").replace("\xa0", " ").strip())


def find(prs, needle):
    """Index of the first slide whose title contains `needle` (case-insensitive)."""
    for i, s in enumerate(prs.slides):
        if needle.lower() in title_of(s).lower():
            return i
    raise SystemExit(f"slide not found by title: {needle!r}")


def set_para(para, text):
    """Replace paragraph text, keeping the first run's formatting."""
    if para.runs:
        para.runs[0].text = text
        for r in para.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        para.text = text


def textbox(slide, left, top, w, h, fill=None, border=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    else:
        box.fill.background()
    if border is not None:
        box.line.color.rgb = border
        box.line.width = Pt(1)
    else:
        box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    # Without this the frame shrinks to its text and the declared height is
    # ignored, which left the coloured panels floating above a band of dead space.
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = Inches(0.12)
    tf.margin_top = tf.margin_bottom = Inches(0.08)
    return box


def write(tf, lines, size=14, bold_first=False, color=INK, space_after=4):
    """lines: list of (text, bold) or plain strings."""
    tf.clear()
    for i, item in enumerate(lines):
        text, bold = item if isinstance(item, tuple) else (item, False)
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = text
        run.font.name = BODY_FONT
        run.font.size = Pt(size)
        run.font.bold = bold or (bold_first and i == 0)
        run.font.color.rgb = color
        para.space_after = Pt(space_after)


def move_slide(prs, from_idx, to_idx):
    lst = prs.slides._sldIdLst
    els = list(lst)
    el = els[from_idx]
    lst.remove(el)
    lst.insert(to_idx, el)


def add_slide_at(prs, layout_name, index):
    layout = {lay.name: lay
              for m in prs.slide_masters for lay in m.slide_layouts}[layout_name]
    slide = prs.slides.add_slide(layout)          # lands at the end
    move_slide(prs, len(prs.slides._sldIdLst) - 1, index)
    return slide


def widen_title(slide, width_in=11.29, size_pt=26):
    """Match the hand-built slides' title box AND type size.

    Two separate traps. The "Title and Content" layout gives a 10.00" title box
    where the existing use-case slides use 11.29"; and the layout's default title
    type is far larger than the 26 pt every hand-built slide uses. Fixing only
    the width still wrapped to two lines, and the overflow rode up out of the
    green band instead of down.
    """
    title = slide.shapes.title
    if title is None:
        return
    # Full geometry copied from the hand-built use-case slide; the layout's box
    # starts at T0.01 and the text rode up off the top of the green band.
    title.left, title.top = Inches(0.76), Inches(0.23)
    title.width, title.height = Inches(width_in), Inches(0.69)
    for para in title.text_frame.paragraphs:
        for r in para.runs:
            r.font.size = Pt(size_pt)
            r.font.name = BODY_FONT


def drop_placeholder(slide, ph_type_name):
    for sh in list(slide.placeholders):
        if str(sh.placeholder_format.type).startswith(ph_type_name):
            sh._element.getparent().remove(sh._element)


# --------------------------------------------------------------------------
# 1. speaker notes
# --------------------------------------------------------------------------

def strip_notes(prs):
    emptied = dropped = 0
    for slide in prs.slides:
        if slide.has_notes_slide:
            slide.notes_slide.notes_text_frame.text = ""
            emptied += 1
        for rId, rel in list(slide.part.rels.items()):
            if rel.reltype == RT.NOTES_SLIDE:
                slide.part.drop_rel(rId)
                dropped += 1
    note(f"[1] notes: emptied {emptied}, detached {dropped}")


# --------------------------------------------------------------------------
# 2-3. font + typo
# --------------------------------------------------------------------------

def normalize_fonts(prs, indices):
    total = 0
    for i in indices:
        n = 0
        for sh in prs.slides[i].shapes:
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    if r.font.name and r.font.name != BODY_FONT:
                        r.font.name = BODY_FONT
                        n += 1
        note(f"[2] slide {i+1} ({title_of(prs.slides[i])[:34]!r}): {n} runs -> {BODY_FONT}")
        total += n
    return total


def fix_typo(prs):
    for s in prs.slides:
        if s.shapes.title is None:
            continue
        for para in s.shapes.title.text_frame.paragraphs:
            for r in para.runs:
                if "Aanalyze" in r.text:
                    r.text = r.text.replace("Aanalyze", "Analyze")
                    note(f"[3] typo fixed: {r.text!r}")


# --------------------------------------------------------------------------
# 4. use case 2 -- lead with the question
# --------------------------------------------------------------------------

def reframe_use_case_2(prs):
    i = find(prs, "AI-Assisted Vaccine Discovery")
    s = prs.slides[i]
    set_para(s.shapes.title.text_frame.paragraphs[0],
             "Use Case: From a Vaccine to Its Host-Gene Network")

    boxes = [sh for sh in s.shapes
             if sh.has_text_frame and "TEXT_BOX" in str(sh.shape_type)]
    boxes.sort(key=lambda sh: sh.top)

    write(boxes[0].text_frame, [
        "A group developing a COVID-19 DNA vaccine asks: which host genes does "
        "the literature already link to this vaccine, and what do they do?",
    ], size=16)

    write(boxes[1].text_frame, [
        "VacNet resolves the vaccine through the Vaccine Ontology and returns its "
        "literature-linked gene network. VacSummarAI summarizes only the retrieved "
        "evidence, so every statement traces back to a PMID.",
    ], size=12)
    note(f"[4] slide {i+1}: reframed as a question-led use case")


# --------------------------------------------------------------------------
# 5. NEW use case 3 -- ontology-guided convergence (R01 preliminary data)
# --------------------------------------------------------------------------

def add_use_case_3(prs, after_title):
    idx = find(prs, after_title) + 1
    s = add_slide_at(prs, "Title and Content", idx)

    s.shapes.title.text_frame.paragraphs[0].text = ""
    # Kept to ~47 characters: the sibling use-case title fits one line at that
    # length, and a second line is clipped by the green band.
    set_para(s.shapes.title.text_frame.paragraphs[0],
             "Use Case: Convergence Across Four Brain Injuries")
    widen_title(s)
    drop_placeholder(s, "BODY")

    q = textbox(s, 0.45, 1.15, 5.30, 0.95, fill=GREEN_BOX, border=TEAL)
    write(q.text_frame, [
        "Do four distinct brain injuries converge on a shared set of "
        "pathogenic pathways?",
    ], size=15)

    m = textbox(s, 0.45, 2.28, 5.30, 3.45)
    write(m.text_frame, [
        ("Ignet supplied one of two independent gene sets.", True),
        "Its Human Disease Ontology layer gave MEDLINE-wide co-occurrence between "
        "HDO / DrugBank identifiers and gene mentions.",
        "A separate PubMed sweep (669,996 PMIDs across 11 disease sets) supplied "
        "the second, independently.",
        "A gene counted only if it appeared in both, at three or more shared PMIDs.",
        ("4 injuries  x  4 pathologies  x  3 comorbidities", True),
    ], size=12, space_after=7)

    if FIGURE.exists():
        # 2.293:1 native. Sized to the right column and dropped to sit between
        # the question box and the result bar rather than floating high.
        s.shapes.add_picture(str(FIGURE), Inches(5.90), Inches(2.05),
                             width=Inches(7.15), height=Inches(3.12))
    else:
        note(f"[5] WARNING figure missing: {FIGURE}")

    r = textbox(s, 0.45, 5.55, 12.45, 1.15, fill=CYAN_BOX)
    write(r.text_frame, [
        ("11 pan-injury genes   |   122 core convergent   |   9 pathways significant "
         "across all 4 injuries AND all 4 outcomes", True),
        "An independent public multi-omics reanalysis later confirmed the "
        "literature-derived signature (Stouffer z = +4.30, p = 1.7 x 10⁻⁵). "
        "Preliminary data for an NIH multi-PI R01 application.",
    ], size=11, space_after=2)

    note(f"[5] inserted new use case at position {idx+1}")
    return idx


# --------------------------------------------------------------------------
# 6. ontology slide -- terms vs structure
# --------------------------------------------------------------------------

def rebuild_ontology_slide(prs):
    old = find(prs, "AI Layer Runs on the Ontology")
    s = add_slide_at(prs, "Title and Content", old)   # insert just before the old one

    set_para(s.shapes.title.text_frame.paragraphs[0],
             "Ontologies in Ignet: Terms vs. Structure")
    widen_title(s)
    drop_placeholder(s, "BODY")

    left = textbox(s, 0.45, 1.22, 6.05, 3.00, fill=RGBColor(0xF4, 0xF6, 0xF7),
                   border=RGBColor(0xC9, 0xD3, 0xD8))
    write(left.text_frame, [
        ("Vocabulary only", True),
        "identifiers stored, hierarchy never traversed",
        "",
        "GO, PSI-MI  —  carried inside the interaction annotation",
        "DrugBank  —  a curated database, used as a term list",
    ], size=14, space_after=8)

    right = textbox(s, 6.83, 1.22, 6.05, 3.00, fill=GREEN_BOX, border=TEAL)
    write(right.text_frame, [
        ("Structure in use", True),
        "the hierarchy does measurable work",
        "",
        "VO  —  subsumption expansion, on by default",
        "INO  —  grouped by class, not by matched phrase",
        "DOID / HDO  —  generality by descendant count, not depth",
    ], size=14, space_after=8)

    bar = textbox(s, 0.45, 4.45, 12.43, 1.85, fill=CYAN_BOX)
    write(bar.text_frame, [
        ("Every AI surface retrieves over this layer.", True),
        "BioBERT scoring, RAG retrieval in the Literature Assistant, GPT-4o synthesis "
        "in BioSummarAI and VacSummarAI, and 8 MCP tools all read the same "
        "ontology-annotated evidence.",
        "So while INO was grouped by matched phrase, every AI consumer was served "
        "“associated / effects / including” as interaction types. Changing how a "
        "class is grouped changes what every downstream model sees.",
    ], size=13, space_after=9)

    # remove the superseded slide (now one position later)
    lst = prs.slides._sldIdLst
    els = list(lst)
    lst.remove(els[old + 1])
    note(f"[6] rebuilt ontology slide at position {old+1}; removed the old one")


# --------------------------------------------------------------------------
# 7. conclusion -- trim
# --------------------------------------------------------------------------

def trim_conclusion(prs):
    i = find(prs, "Conclusion")
    s = prs.slides[i]

    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text
        if "Final Takeaway" in txt:
            paras = [p for p in sh.text_frame.paragraphs if p.text.strip()]
            set_para(paras[0], "Final Takeaway")
            set_para(paras[1],
                     "Literature mining, biomedical ontologies and AI-assisted ranking, "
                     "combined into one transparent and explainable discovery platform.")
        elif "Strengths" in txt:
            new = [
                ("Strengths", True),
                "Genes, vaccines, diseases and drugs unified through biomedical ontologies",
                "Every interaction traceable to a BioBERT-scored sentence and its PMID",
                "Daily updates; web tools, REST API and MCP",
                ("Limitations", True),
                "PubMed abstracts only — no trials, patents or full text",
                "Confidence reflects literature evidence, not experimental validation",
                "No multi-omics or experimental data yet",
            ]
            paras = [p for p in sh.text_frame.paragraphs if p.text.strip()]
            if len(paras) != len(new):
                note(f"[7] WARNING conclusion has {len(paras)} paragraphs, "
                     f"expected {len(new)} — skipping to avoid mangling")
                continue
            for para, item in zip(paras, new):
                set_para(para, item[0] if isinstance(item, tuple) else item)
    note(f"[7] slide {i+1}: conclusion trimmed")


# --------------------------------------------------------------------------
# 8. benchmark -> backup
# --------------------------------------------------------------------------

def benchmark_to_backup(prs):
    i = find(prs, "Platform Benchmarking")
    last = len(prs.slides._sldIdLst) - 1
    move_slide(prs, i, last)
    note(f"[8] moved Platform Benchmarking from {i+1} to {last+1} (backup)")


# --------------------------------------------------------------------------

def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--in", dest="src", required=True)
    ap.add_argument("--out", dest="dst", required=True)
    a = ap.parse_args()

    src, dst = pathlib.Path(a.src), pathlib.Path(a.dst)
    shutil.copy2(src, dst)
    prs = Presentation(str(dst))
    before = len(prs.slides)

    strip_notes(prs)
    uc1 = find(prs, "Open API + MCP")
    uc2 = find(prs, "AI-Assisted Vaccine Discovery")
    normalize_fonts(prs, [uc1, uc2])
    fix_typo(prs)
    reframe_use_case_2(prs)
    add_use_case_3(prs, "From a Vaccine to Its Host-Gene Network")
    rebuild_ontology_slide(prs)
    trim_conclusion(prs)
    # The conclusion mixed fonts within one slide: the takeaway box was Times
    # New Roman while the body beside it was Trebuchet.
    normalize_fonts(prs, [find(prs, "Conclusion")])
    benchmark_to_backup(prs)

    prs.save(str(dst))
    after = len(Presentation(str(dst)).slides)
    note(f"\nslides: {before} -> {after}   written: {dst}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
