#!/usr/bin/env python3
"""Build the ICIBM 2026 Ontology-Applications deck from Sayed's VDOS deck.

Approach: start FROM the VDOS file rather than a blank presentation, so the
theme, masters, figures and screenshots survive untouched. Then curate the slide
list, retitle for this talk, and insert the ontology-depth material that the
VDOS deck predates.

Why not build from blank: the existing decks carry screenshots and diagrams that
would be lost, and rebuilding the brand theme by hand is exactly the trap in
concept/python-pptx-template-ships-sample-slides. Copy-and-curate keeps every
existing asset and touches only what changes.

The new content is this session's ontology-depth work, which is the part
specific to an ontology-applications audience:
  - ontology TERMS vs ontology STRUCTURE (the audit)
  - INO: matched phrases -> ontology classes (before/after)
  - annotation is already multi-ontology (INO + MI + GO)
  - DOID: classifier filtering by structure, not by a name list
  - VO: class-level subsumption already running in Vignet
  - data quality: obsolete classes, versioned annotation

Usage:
    python3 scripts/build_icibm_deck.py \
        --source presentations/ICIBM2026/VDOS_2026_Workshop_Ignet2_Vignet_FullPaper_Slide_SA_HR_v2.1.pptx \
        --out    presentations/ICIBM2026/ICIBM2026_Ignet2_Vignet_Ontology_JHur_v1.pptx
"""
from __future__ import annotations

import argparse
import pathlib
import re

from PIL import Image
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

# Speaker notes live in a tracked sibling file, NOT in the .pptx: this script
# regenerates the deck from Sayed's source every run, so notes typed into
# PowerPoint would be discarded on the next build. presentations/ is gitignored,
# so the file sits in scripts/ alongside this builder.
NOTES_FILE = pathlib.Path(__file__).with_name("icibm2026_speaker_notes.md")
NOTES_HEADING = re.compile(r"^##\s+(\d+)\s*\|\s*(.+?)\s*$", re.MULTILINE)

# Content layout on master 2. idx 0 = title, idx 15 = bulleted body.
BULLET_LAYOUT = (2, 3)

# VDOS slides to carry over, in source order, with why each earns its place.
# (1-based source index, short note)
#
# Curated for a 20-MINUTE ICIBM slot (~15 min speaking + 5 min Q&A). At ~60 s
# per slide that budget buys ~19 slides, so seven VDOS slides come out. Four are
# genuine cuts; three are figure/prose pairs where the prose slide was really
# speaker-note text and moves into the notes of the figure it described.
KEEP = [
    (1,  "title - rewritten below"),
    # (2)  talk roadmap        - CUT: a roadmap slide does not survive 15 min.
    (3,  "literature mining challenges + gap"),
    # (4)  objectives          - CUT: redundant once slide 3 states the gap and
    #                            'Ontology Depth' states this talk's thesis.
    (5,  "dual platform: Ignet 2.0 + Vignet"),
    (8,  "integrated ontologies - amended below"),
    (9,  "semantic database architecture"),
    (10, "Ignet 2.0 platform overview (figure)"),
    # (11) Ignet discovery     - MERGED into 10's notes: prose describing the
    #                            figure, plus a stale 7/7/2026 data stamp.
    (12, "Vignet platform overview (figure)"),
    # (13) Vignet exploration  - MERGED into 12's notes, same reason as 11.
    # (14) programmatic access - MERGED into 16: REST + MCP facts belong on the
    #                            slide that shows MCP actually being used.
    (16, "open API + MCP - use case 1, amended below"),
    (17, "use case: AI-assisted vaccine discovery"),
    (18, "platform benchmarking"),
    (19, "conclusion"),
    # (20) future directions   - CUT: 'Ontology Hygiene' already closes forward.
    (21, "acknowledgement"),
    (22, "availability - now the closing slide, stays up through Q&A"),
    # (23) questions           - CUT: 22 carries the QR code and URLs, which are
    #                            more useful on screen during Q&A than 'Any
    #                            Question?' is.
]

# New slides, each inserted AFTER the given source slide number once the deck
# has been curated. Body lines: (text, indent_level).
NEW_SLIDES = [
    {
        "after": 8,
        "title": "Ontology Depth: Terms vs. Structure",
        "body": [
            ("The question is not WHICH ontologies are integrated, but whether "
             "their STRUCTURE is used.", 0),
            ("VO — structural: is_a hierarchy drives class-level queries", 0),
            ("INO — was vocabulary only: identifiers stored but unused", 0),
            ("DOID — was vocabulary only: diseases as free text", 0),
            ("This talk: what changed when we used the structure "
             "we already had.", 0),
            ("And why that matters MORE, not less, once an AI agent is the "
             "consumer.", 0),
        ],
    },
    {
        "after": 8,
        "title": "Annotation Is Already Multi-Ontology",
        "body": [
            ("The 166 interaction classes in use span three ontologies:", 0),
            ("86 INO  ·  57 MI (PSI Molecular Interactions)  ·  23 GO", 1),
            ("Plus 15 INO_T* template artifacts — not ontology classes, "
             "excluded by namespace rather than by a stopword list.", 0),
            ("Interaction typing is already an ontology ENSEMBLE; naming only "
             "INO understates the semantic layer in place.", 0),
        ],
    },
    {
        "after": 9,
        "title": "INO: Phrases to Ontology Classes",
        "body": [
            ("Every annotation already carried an INO identifier (54.7M rows), "
             "but views grouped on the matched PHRASE.", 0),
            ("INO_0000157 alone spans 69 phrases: 'effects', 'effect', 'changes'", 1),
            ("BEFORE — top 'interaction types':", 0),
            ("associated · effects · induced · expression · including", 1),
            # The screenshot IS the "after", with live counts. Listing them here
            # too duplicated the image and pushed the punchline off the slide.
            ("AFTER — by ontology class: the live panel on the right.", 0),
            ("Same data. The semantics were already in the identifiers.", 0),
        ],
    },
    {
        "after": 9,
        "title": "DOID: Classifiers by Structure",
        "body": [
            ("Over-broad disease terms were removed by a hand-written list: "
             "{ disease, syndrome, disorder }", 0),
            ("'disease' is DOID:4 — the ROOT — filtered by its label. The list "
             "missed 'cancer' (1.3M) and 'carcinoma' (244k).", 0),
            ("Depth fails too: 'fibromyalgia' is level 2, same as 'cancer'.", 0),
            ("Transitive descendant count separates them:", 0),
            ("disease 12,220 · cancer 2,160 · carcinoma 622  ||  "
             "breast cancer 86 · fibromyalgia 0", 1),
            ("Structure removes 4.70M classifier annotations vs 3.07M by name, "
             "dropping no real diagnosis.", 0),
        ],
    },
    {
        # Anchored on 12 (Vignet platform overview), not the retired 13.
        "after": 12,
        "title": "VO: Class-Level Reasoning in Vignet",
        "body": [
            ("Selecting a vaccine class retrieves evidence for its is_a "
             "descendants — subsumption, on by default.", 0),
            # Says what the audience can actually SEE. The earlier wording
            # ("aggregates 239 data-bearing descendants") named a number that
            # appears nowhere on screen -- the rendered network carries exactly
            # ONE vaccine node; the descendants' evidence rolls up rather than
            # rendering. The node/edge counts ARE on screen, so anchor to those.
            ("'viral vaccine' rolls up 239 data-bearing descendant classes: "
             "90 -> 1,021 nodes, 89 -> 3,249 edges", 1),
            ("6,796 VO classes, subtree-closed, so parent classes stay "
             "meaningful selections rather than dead ends.", 0),
            # The closing line ("a question asked at the level the researcher
            # thinks in...") is a spoken flourish and overflowed behind the
            # screenshots. It lives in the speaker notes instead.
        ],
    },
    {
        # Also anchored on 12, declared AFTER the VO slide so it lands directly
        # after it: introduce the AI stack, then the two use cases show it running.
        "after": 12,
        "title": "The AI Layer Runs on the Ontology Layer",
        "body": [
            ("BioBERT — a deployed model service; every mined interaction "
             "sentence carries its confidence score.", 0),
            ("RAG — the Literature Assistant retrieves our own evidence "
             "sentences and grounds the answer in cited PMIDs.", 0),
            ("GPT-4o — BioSummarAI and VacSummarAI synthesize retrieved "
             "evidence, not free-form recall.", 0),
            ("MCP — 8 tools, so an external agent queries the same "
             "ontology-guided knowledge base.", 0),
            ("Every one of these RETRIEVES over the ontology layer. While INO "
             "grouped by phrase, every AI consumer was served 'associated / "
             "effects / including' as interaction types.", 0),
        ],
    },
    {
        "after": 18,
        "title": "Ontology Hygiene",
        "body": [
            ("Using structure surfaces quality signals that flat term matching "
             "hides:", 0),
            ("6 of 23 referenced GO classes are OBSOLETE", 1),
            ("2,514 of 14,735 DOID terms are obsolete", 1),
            ("An identifier is not self-validating — it must resolve against a "
             "VERSIONED ontology release to be trusted.", 0),
            ("Annotation pipelines should record the ontology version used, "
             "and re-resolve on refresh.", 0),
        ],
    },
]

# Screenshots placed onto new slides, keyed by slide title. Paths are relative
# to the --source deck's directory, which is gitignored -- the same directory the
# source .pptx lives in, so requiring them here adds no new precondition.
#
# `body` overrides the bullet placeholder's box (inches) to make room; `images`
# are aspect-fit and centred inside their own boxes, never stretched.
IN = 914400  # EMU per inch

SLIDE_IMAGES = {
    "INO: Phrases to Ontology Classes": {
        "body": (0.60, 1.26, 5.70, 5.40),
        "images": [
            # The live Explorer IS the "after" evidence; text sits left of it.
            {"file": "screenshots/crop-slide08-ino.png",
             "box": (6.55, 1.40, 6.30, 5.20)},
        ],
    },
    "VO: Class-Level Reasoning in Vignet": {
        "body": (1.25, 1.26, 10.83, 1.95),
        "images": [
            {"file": "screenshots/crop-slide12-off.png",
             "box": (0.90, 3.40, 5.30, 3.30),
             "caption": "subsumption OFF — 90 nodes, 89 edges"},
            {"file": "screenshots/crop-slide12-on.png",
             "box": (6.70, 3.40, 5.30, 3.30),
             "caption": "subsumption ON — 1,021 nodes, 3,249 edges"},
        ],
    },
}

TITLE_MAIN = ("From Biomedical Ontologies to Knowledge Discovery:\n"
              "Ontology-Driven Applications in Ignet 2.0 and Vignet")
TITLE_SUB = (
    "Junguk Hur, PhD\n"
    "Associate Professor, Department of Biomedical Sciences\n"
    "University of North Dakota\n"
    "\n"
    "ICIBM 2026 — Biological and Biomedical Ontology Workshop in the AI Era"
)


def add_bullet_slide(prs, title: str, body: list[tuple[str, int]]):
    master = list(prs.slide_masters)[BULLET_LAYOUT[0]]
    layout = master.slide_layouts[BULLET_LAYOUT[1]]
    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = title
    # The brand band fits ONE line; LibreOffice render showed 2-line titles
    # clipped at the band edge. Cap the size and disable autofit growth.
    title_tf = slide.shapes.title.text_frame
    title_tf.word_wrap = True
    for para in title_tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(26)

    body_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 15:
            body_ph = ph
            break
    if body_ph is None:  # layout drift: fall back to any BODY placeholder
        body_ph = next(
            ph for ph in slide.placeholders if ph.placeholder_format.idx != 0
        )

    tf = body_ph.text_frame
    tf.word_wrap = True  # python-pptx re-centres without this
    tf.clear()
    for i, (text, level) in enumerate(body):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = text
        para.level = level
        size = 16 if level == 0 else 14 if level == 1 else 12
        for run in para.runs:
            run.font.size = Pt(size)
    return slide


def place_images(slide, spec: dict, asset_root: pathlib.Path) -> None:
    """Resize the bullet box and add aspect-fit screenshots beside/below it.

    Aspect-fit, never stretch: scale by whichever of width/height binds first,
    then centre the picture inside its box. A stretched screenshot of a network
    graph reads as a rendering fault, not a design choice.

    Missing assets raise rather than skip. presentations/ is gitignored, so an
    absent file means the checkout is incomplete -- and a slide that silently
    loses its evidence is the failure mode this builder has already hit twice.
    """
    if "body" in spec:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 15:
                bx, by, bw, bh = spec["body"]
                ph.left, ph.top, ph.width, ph.height = (
                    Emu(int(bx * IN)), Emu(int(by * IN)),
                    Emu(int(bw * IN)), Emu(int(bh * IN)),
                )
                break

    for img in spec.get("images", []):
        path = asset_root / img["file"]
        if not path.is_file():
            raise SystemExit(
                f"missing slide asset: {path}\n"
                f"  (regenerate the crops, or drop the entry from SLIDE_IMAGES)"
            )
        bl, bt, bw, bh = img["box"]
        with Image.open(path) as im:
            iw, ih = im.size
        scale = min(bw / (iw / 96), bh / (ih / 96))  # PNGs are 96 dpi here
        w_in, h_in = (iw / 96) * scale, (ih / 96) * scale
        left = bl + (bw - w_in) / 2
        top = bt + (bh - h_in) / 2
        slide.shapes.add_picture(
            str(path), Emu(int(left * IN)), Emu(int(top * IN)),
            Emu(int(w_in * IN)), Emu(int(h_in * IN)),
        )

        if img.get("caption"):
            box = slide.shapes.add_textbox(
                Emu(int(bl * IN)), Emu(int((bt + bh + 0.06) * IN)),
                Emu(int(bw * IN)), Emu(int(0.32 * IN)),
            )
            tf = box.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            para.text = img["caption"]
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.size = Pt(13)
                run.font.bold = True


def parse_notes(path: pathlib.Path) -> dict[int, str]:
    """Parse the speaker-notes markdown into {slide_position: notes_text}.

    Sections are `## <position> | <title>`; the body runs to the next heading.
    Position is the slide's place in the FINAL deck, so this must be applied
    after the reorder, not against source indices.
    """
    text = path.read_text(encoding="utf-8")
    matches = list(NOTES_HEADING.finditer(text))
    if not matches:
        raise SystemExit(f"{path}: no '## <n> | <title>' sections found")

    notes: dict[int, tuple[str, str]] = {}
    for i, m in enumerate(matches):
        pos = int(m.group(1))
        if pos in notes:
            raise SystemExit(f"{path}: duplicate section for slide {pos}")
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        body = text[m.end():end].strip()
        if not body:
            raise SystemExit(f"{path}: slide {pos} section is empty")
        notes[pos] = (m.group(2), body)
    return notes


def _norm(s: str) -> str:
    """Lowercase, drop punctuation, collapse whitespace -- for title matching."""
    return " ".join(re.sub(r"[^a-z0-9 ]+", " ", s.lower()).split())


def apply_notes(prs, notes: dict[int, tuple[str, str]]) -> None:
    """Replace every slide's notes with ours.

    Replace, never append: the inherited VDOS notes are in a different speaker's
    voice (slide 1's opened "My name is Sayed Asaduzzaman") and are paced for a
    much longer slot, so leaving any of them in place would be worse than having
    none. Every slide must be covered -- a silent gap would surface as someone
    else's script in presenter view mid-talk.

    Two checks, because they fail differently. Coverage catches a MISSING
    section. The title match catches a MISPLACED one -- notes attached to the
    wrong slide, which coverage cannot see because every position is still
    filled. Inserting a slide renumbers everything after it, so that is the
    likely failure whenever the deck's shape changes (observed once already,
    when a renumber swapped Use Case 2 with Platform Benchmarking).
    """
    missing = [i for i in range(1, len(prs.slides) + 1) if i not in notes]
    if missing:
        raise SystemExit(
            f"{NOTES_FILE.name}: no notes for slide(s) {missing}; "
            f"deck has {len(prs.slides)} slides"
        )
    extra = [i for i in notes if i > len(prs.slides)]
    if extra:
        raise SystemExit(
            f"{NOTES_FILE.name}: notes for slide(s) {extra} but deck has "
            f"only {len(prs.slides)} slides -- did the trim change?"
        )

    for pos, (heading, body) in sorted(notes.items()):
        slide = prs.slides[pos - 1]
        actual = slide.shapes.title.text if slide.shapes.title is not None else ""
        want, got = _norm(heading), _norm(actual)
        # Slide 1's title runs onto a second line, so prefix-match rather than
        # require equality.
        if not (got == want or got.startswith(want) or want.startswith(got)):
            raise SystemExit(
                f"{NOTES_FILE.name}: slide {pos} mismatch -- notes say "
                f"{heading!r} but the slide is {actual!r}. Notes are keyed by "
                f"FINAL position; did a slide get inserted or removed?"
            )

        tf = slide.notes_slide.notes_text_frame
        tf.clear()
        for i, line in enumerate(body.split("\n")):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = line


def normalize_kept_titles(prs, keep_indices: list[int]) -> None:
    """Cap inherited slide titles so they cannot wrap into the brand band.

    add_bullet_slide already caps titles on the slides this script creates,
    because a 2-line title renders clipped at the band edge. The inherited VDOS
    slides were never covered by that, and a full-deck render caught two of them
    clipped: 'The Literature Mining Challenges and Gap' and 'Use Case 2:
    AI-Assisted Vaccine Knowledge Discovery'.

    Capping is defensive: PowerPoint's own autofit may shrink these anyway, but
    a smaller title never renders worse, and the LibreOffice render is the only
    check available here.
    """
    for idx in keep_indices:
        slide = prs.slides[idx - 1]
        if slide.shapes.title is None:
            continue
        tf = slide.shapes.title.text_frame
        tf.word_wrap = True
        for para in tf.paragraphs:
            for run in para.runs:
                if run.font.size is None or run.font.size > Pt(26):
                    run.font.size = Pt(26)


def shorten_use_case_2_title(prs) -> None:
    """Slide 17's title is long enough to wrap even at the 26 pt cap.

    'Knowledge' is redundant between 'Vaccine' and 'Discovery', and dropping it
    takes the title from 51 to 41 characters -- comfortably inside one line at
    the same size the neighbouring MCP slide renders at.
    """
    slide = prs.slides[16]  # 0-based; source slide 17
    if slide.shapes.title is None:
        raise RuntimeError("slide 17: no title placeholder to shorten")
    tf = slide.shapes.title.text_frame
    tf.word_wrap = True
    tf.text = "Use Case 2: AI-Assisted Vaccine Discovery"
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(26)


def amend_mcp_slide(prs) -> None:
    """Fold retired slide 14 (programmatic access) into slide 16 (MCP use case).

    Slide 14 was a text-only slide stating that both platforms expose a REST API
    and an 8-tool MCP server. Slide 16 SHOWS that access being used, so 16 is
    retitled to own both claims.

    Retitle ONLY -- deliberately no body text. The first attempt appended those
    facts as a third line in 16's explanatory box; the box is sized for two
    bullets and the render showed the line crossing its bottom border in a
    smaller unbulleted font. Beyond the overflow, a dense one-liner earns
    nothing on a slide the audience sees for ~40 s: 'REST API + 8 MCP tools' is
    a spoken point, and it lives in this slide's speaker notes instead.
    """
    slide = prs.slides[15]  # 0-based; source slide 16
    if slide.shapes.title is None:
        raise RuntimeError("slide 16: no title placeholder to retitle")

    tf = slide.shapes.title.text_frame
    tf.word_wrap = True
    tf.text = "Open API + MCP: LLM-Powered Knowledge Discovery"
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(26)


def set_title_slide(prs) -> None:
    """Rewrite the title slide for this talk.

    The source title slide carries THREE text frames: headline, author block and
    a superscript affiliations list. Setting only the first two leaves the third
    in place, and the render showed the new subtitle drawn straight over it. Any
    frame beyond the two we own is emptied.
    """
    slide = prs.slides[0]
    frames = [sh for sh in slide.shapes
              if sh.has_text_frame and sh.text_frame.text.strip()]
    if not frames:
        return
    frames.sort(key=lambda sh: sh.top or 0)

    frames[0].text_frame.word_wrap = True
    frames[0].text_frame.text = TITLE_MAIN
    for para in frames[0].text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(36)
            run.font.name = BODY_FONT

    if len(frames) > 1:
        tf = frames[1].text_frame
        tf.word_wrap = True
        tf.text = TITLE_SUB
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(20)
                run.font.name = BODY_FONT

    # Clear every remaining frame so nothing from the source deck shows through.
    for extra in frames[2:]:
        extra.text_frame.clear()


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--source", required=True)
    ap.add_argument("--out", required=True)
    args = ap.parse_args()

    prs = Presentation(args.source)
    original_count = len(prs.slides)

    # ORDER MATTERS: add new slides BEFORE dropping any.
    # add_slide derives its partname from the current slide count, so dropping
    # first frees numbers that retained slides still occupy -- the new parts then
    # collide (observed: duplicate ppt/slides/slide21..23.xml in the saved zip,
    # which PowerPoint reports as a file needing repair). Appending while the
    # deck is still full-length forces fresh partnames.
    keep_indices = [idx for idx, _ in KEEP]

    # A new slide is placed by appending it after its anchor as the kept slides
    # are walked, so an anchor that is not itself kept never gets visited: the
    # new slide is silently omitted from `order` and its part orphaned, with no
    # error. That is a whole slide disappearing quietly, so fail loudly instead.
    orphaned = sorted({spec["after"] for spec in NEW_SLIDES} - set(keep_indices))
    if orphaned:
        raise SystemExit(
            f"anchor slide(s) {orphaned} are in NEW_SLIDES but not in KEEP; "
            f"re-anchor those new slides onto a kept source slide"
        )

    added: list[tuple[int, int]] = []  # (source anchor, position appended at)
    asset_root = pathlib.Path(args.source).resolve().parent
    for spec in NEW_SLIDES:
        slide = add_bullet_slide(prs, spec["title"], spec["body"])
        if spec["title"] in SLIDE_IMAGES:
            place_images(slide, SLIDE_IMAGES[spec["title"]], asset_root)
        added.append((spec["after"], len(prs.slides) - 1))

    # Amend kept slides while source indices are still their original ones.
    normalize_kept_titles(prs, keep_indices)
    shorten_use_case_2_title(prs)
    amend_mcp_slide(prs)

    # Desired final order, expressed as positions in the CURRENT deck.
    order: list[int] = []
    for src in keep_indices:
        order.append(src - 1)
        for anchor, pos in added:
            if anchor == src:
                order.append(pos)

    # Rebuild sldIdLst in that order; anything absent is dropped along with its
    # relationship, so no orphan parts are left behind.
    id_list = prs.slides._sldIdLst
    entries = list(id_list)
    keep_set = set(order)
    for entry in entries:
        id_list.remove(entry)
    for pos in order:
        id_list.append(entries[pos])
    for pos, entry in enumerate(entries):
        if pos not in keep_set:
            prs.part.drop_rel(entry.rId)

    set_title_slide(prs)

    # After the reorder: notes are keyed by FINAL slide position.
    notes = parse_notes(NOTES_FILE)
    apply_notes(prs, notes)

    prs.save(args.out)

    print(f"source slides : {original_count}")
    print(f"kept          : {len(keep_indices)}")
    print(f"new slides    : {len(NEW_SLIDES)}")
    print(f"final         : {len(prs.slides)}")
    print(f"notes applied : {len(notes)} (from {NOTES_FILE.name})")
    print(f"written       : {args.out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
